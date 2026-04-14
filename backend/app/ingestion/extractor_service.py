import fitz  # PyMuPDF
from pptx import Presentation
import os
import re

# Patterns for academic paper noise (appendices, figure OCR, citations)
_ARXIV_PATTERN = re.compile(
    r"arXiv:\d+\.\d+(?:v\d+)?\s*\[[\w.]+\]\s*\d+\s+\w+\s+\d{4}",
    re.IGNORECASE
)
# Common QA/prompt template leakage from paper appendices
_QA_PROMPT_PATTERNS = [
    re.compile(r"You have access to memories from[^.]+\.", re.IGNORECASE | re.DOTALL),
    re.compile(r"The answer should be (?:less than|at most) \d+[-–]?\d*\s*words?", re.IGNORECASE),
    re.compile(r"Answer the following question based on[^.]+\.", re.IGNORECASE),
    re.compile(r"\.\s*You have \d+ memories[^.]+\.", re.IGNORECASE | re.DOTALL),
]
# Reference-like lines (author et al., year - common in references section)
_REF_LINE = re.compile(
    r"^[\w\s,.-]+\set\s+al\.\s*[,.]?\s*\d{4}",
    re.IGNORECASE
)
# Figure/Table caption noise (short lines that are likely OCR from figures)
_CAPTION_NOISE = re.compile(
    r"^(?:Figure|Table|Fig\.|Tab\.)\s*\d*[.:]?\s*$",
    re.IGNORECASE
)
# Additional junk tags seen in noisy extractions
_JUNK_PATTERNS = [
    re.compile(r"Back to (?:the page|Mail Online|top)", re.IGNORECASE),
    re.compile(r"Slide \d+ of \d+", re.IGNORECASE),
    re.compile(r"Page \d+ of \d+", re.IGNORECASE),
    re.compile(r"\[\d+\]", re.IGNORECASE), # Footnote citations as junk when standalone
    re.compile(r"^\s*[\W_]+\s*$", re.IGNORECASE), # Purely non-alphanumeric lines
]

# Emails and "First Last, Ph.D (…)" lines in pasted decks / model echoes
_EMAIL_ANYWHERE = re.compile(
    r"\(?\s*[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\s*\)?"
)
_NAME_PHD_ATTRIB = re.compile(
    r"\b(?:[A-Z][a-z]+\s+){1,3}[A-Z][a-z]+,\s*Ph\.D\s*(?:\([^)]{0,120}\))?\s*",
    re.UNICODE,
)
_URL_INLINE = re.compile(r"https?://\S+|www\.\S+", re.IGNORECASE)


class ExtractorService:
    @staticmethod
    def normalize_pipeline_text(text: str) -> str:
        """
        Remove lecturer attribution and emails from pasted or extracted text (and from summaries).
        Keeps course content; replaces name+Ph.D with neutral 'The instructor'.
        """
        if not text or not str(text).strip():
            return text
        t = str(text)
        t = _EMAIL_ANYWHERE.sub(" ", t)
        t = _NAME_PHD_ATTRIB.sub("The instructor ", t)
        t = _URL_INLINE.sub(" ", t)
        t = re.sub(r"\*?\s*Wikipedia\*?", " ", t, flags=re.IGNORECASE)
        t = re.sub(r"(?:\s*The instructor\s*){2,}", " The instructor ", t, flags=re.IGNORECASE)
        t = re.sub(r"[ \t]+", " ", t)
        t = re.sub(r"\s+\.\s+", ". ", t)
        t = re.sub(r"\s*\.\s*\.", ".", t)
        return t.strip()

    @staticmethod
    def extract_text(file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == '.pdf':
                text = ExtractorService._extract_from_pdf(file_path)
            elif ext == '.pptx':
                text = ExtractorService._extract_from_pptx(file_path)
            elif ext == '.txt':
                text = ExtractorService._extract_from_txt(file_path)
            else:
                raise ValueError(f"Unsupported file extension: {ext}")
            
            # 1. Strip headers first before ANY sanitization to ensure patterns match accurately
            text = ExtractorService._strip_headers(text)
            
            # 2. Basic sanitization
            text = ExtractorService._sanitize_text(text)
            
            # 3. Extra cleanup specifically for PDFs (e.g. arXiv tags, OCR noise)
            if ext == '.pdf':
                text = ExtractorService._clean_academic_noise(text)
            return text
        except Exception as e:
            raise RuntimeError(f"Failed to extract text from {file_path}: {str(e)}")

    @staticmethod
    def _strip_headers(text: str) -> str:
        """Strip Title, Authors, Emails, Affiliations from the beginning of a document."""
        if not text or len(text.strip()) < 100:
            return text

        lines = text.splitlines()
        content_start_idx = 0
        header_patterns = [
            re.compile(r"\{?[\w.-]+@[\w.-]+\}?"), # Emails
            re.compile(r"(?:University|Institute|Department|Laboratory|School|College|Academy|Faculty of)", re.IGNORECASE), # Affiliations
            re.compile(r"^\d+(?:st|nd|rd|th)?\s+(?:Author|Writer|Contributor)$", re.IGNORECASE),
            re.compile(r"^\*?\s*Corresponding\s+author", re.IGNORECASE),
            re.compile(r"DSC\d{4}:?", re.IGNORECASE), # Course codes
            re.compile(r"Lecture\s+\d+", re.IGNORECASE),
            re.compile(r"Topic:\s+", re.IGNORECASE), # Topic headers
            re.compile(r"Simon\s+Fred", re.IGNORECASE), # Specific author names
            re.compile(r"Department\s+of", re.IGNORECASE)
        ]
        
        # Look at the first 30 lines
        for i, line in enumerate(lines[:30]):
            line = line.strip()
            if not line: continue
            # If we see an Abstract or Introduction, the header definitely ended
            if re.match(r"^(?:Abstract|1\.?\s*Introduction|Introduction|Lecture\s+Objectives|Lecture\s+Overview)", line, re.IGNORECASE):
                content_start_idx = i
                break
            # If line is highly likely metadata, keep looking
            found_match = False
            for p in header_patterns:
                if p.search(line):
                    content_start_idx = i + 1
                    found_match = True
                    break
        
        # Strip the identified header if it's not the whole document
        if content_start_idx > 0 and content_start_idx < len(lines) - 2:
            return "\n".join(lines[content_start_idx:])
        return text

    @staticmethod
    def _clean_academic_noise(text: str) -> str:
        """Remove common academic paper noise: arXiv IDs, QA prompts, figure OCR, reference blocks."""
        if not text or len(text.strip()) < 100:
            return text

        # 1. Remove arXiv-style citation lines
        text = _ARXIV_PATTERN.sub("", text)

        # 2. Split into paragraphs and filter
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        filtered = []
        for p in paragraphs:
            words = p.split()
            # Skip very short fragments (likely figure captions, labels)
            if len(words) < 12 and (_CAPTION_NOISE.match(p) or any(
                x in p.lower() for x in ("figure", "table", "fig.", "tab.")
            )):
                continue
            # Skip very short lines that look like reference entries (avoid body text)
            if len(words) < 12 and _REF_LINE.match(p):
                continue
            filtered.append(p)

        text = "\n\n".join(filtered)

        # 3. Remove stray URLs
        text = re.sub(r"https?://\S+", "", text)
        # Collapse excessive newlines
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    @staticmethod
    def _sanitize_text(text: str) -> str:
        """Sanitize text by removing excessive whitespace, merging hyphens, and truncating references."""
        if not text:
            return ""
        
        # 1. Truncate at References/Bibliography (common in academic papers)
        ref_markers = [
            r'\n\s*References\s*\n', 
            r'\n\s*BIBLIOGRAPHY\s*\n', 
            r'\n\s*Works Cited\s*\n'
        ]
        for marker in ref_markers:
            match = re.search(marker, text, re.IGNORECASE)
            if match:
                text = text[:match.start()]
                break

        # 2. Fix hyphenated line-breaks: "syn-\nthesized" -> "synthesized"
        text = re.sub(r'(\w+)-\n\s*(\w+)', r'\1\2', text)

        # 3. Apply junk patterns
        for pattern in _JUNK_PATTERNS:
            text = pattern.sub("", text)

        # 4. Strip URLs
        text = re.sub(r'https?://\S+|www\.\S+', '', text)
        
        # 5. Split into lines to identify repeating headers/footers
        lines = text.splitlines()
        # Filter purely decorative lines or common noisy fragments
        lines = [l.strip() for l in lines if len(l.strip()) > 2 or l.strip().isalnum()]
        
        if len(lines) > 20:
            from collections import Counter
            line_counts = Counter(lines)
            # Remove any line that appears in more than 10% of the document (likely a header/footer)
            # Tightened from 15% to 10%
            threshold = max(2, len(lines) // 10) 
            text = "\n".join([line for line in lines if line_counts[line] < threshold])
        else:
            text = "\n".join(lines)

        # Remove control characters
        text = "".join(char for char in text if char.isprintable() or char in "\n\r\t")
        
        # Normalize whitespace 
        text = re.sub(r' +', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text) # Allow double newlines but not triple+
        return text.strip()

    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint file, handling GroupShape and other complex shapes."""
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                # Handle different shape types - GroupShape doesn't have 'text' attribute
                try:
                    if hasattr(shape, "text") and shape.text:
                        text += str(shape.text) + "\n"
                except (AttributeError, TypeError):
                    # Skip shapes that can't be read (e.g., GroupShape, embedded objects)
                    pass
        return text

    @staticmethod
    def _extract_from_txt(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

